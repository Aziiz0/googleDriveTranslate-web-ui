import re
from deep_translator import GoogleTranslator
from string import punctuation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from pptx import Presentation
from docx import Document
import urllib.parse
import io
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.http import MediaFileUpload
from google.oauth2 import service_account
from googleapiclient.discovery import build
from google_auth_oauthlib.flow import InstalledAppFlow
import time
import win32com.client
import pythoncom

start_translating = False

# Initialize translator
#translator = Translator()

# Set up OAuth 2.0 credentials
client_id = '436858784684-oef4d0qhfegmn7e0fkikbes0vemk6l01.apps.googleusercontent.com'
client_secret = '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'
scopes = ['https://www.googleapis.com/auth/drive']
redirect_uri = 'http://localhost:58936/'

# Set up token file
token_file = 'token.json'

# Authenticate and authorize the user
flow = InstalledAppFlow.from_client_secrets_file(
    client_secrets_file=os.path.abspath('client_secrets.json'),
    scopes=scopes,
    redirect_uri=redirect_uri
)
auth_url, _ = flow.authorization_url(prompt='consent')

# Print the authorization URL
print('Please visit this URL to authorize the application:', auth_url)

# Once authorized, enter the authorization code
authorization_code = input('Enter the authorization code: ')

# Fetch the access token
flow.fetch_token(
    authorization_response=authorization_code,
)

# Save the credentials to a file
credentials_json = flow.credentials.to_json()
with open(token_file, 'w') as token:
    token.write(credentials_json)

# Build the Google Drive API service
drive_service = build('drive', 'v3', credentials=flow.credentials)

def remove_illegal_chars(filename):
    illegal_chars = r"[']"  # Define the pattern for illegal characters (single quote in this case)
    return re.sub(illegal_chars, "", filename)  # Use regular expression to remove illegal characters

def is_punctuation(text):
    return all(char in punctuation for char in text)

def split_text_into_chunks(text, chunk_size=5000):
    """
    Function to split the text into smaller chunks
    """
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def translate_text(text):
    """
    Function to translate the text
    """
    text = str(text)  # Ensure the text is a string

    # Return the original text if it has less than 2 non-whitespace characters 
    # or consists only of punctuation or if it is a digit
    if len(text.strip()) < 2 or is_punctuation(text.strip()) or text.isdigit():
        return text

    translator = GoogleTranslator(source='zh-CN', target='en')  # Create a translator object
    chunks = split_text_into_chunks(text)  # Split the text into chunks
    translated_text = ""  # Placeholder for the translated text

    for chunk in chunks:
        while True:
            try:
                translated_chunk = translator.translate(chunk)  # Translate the chunk
                # print(f"Translated chunk: {translated_chunk}")
                translated_text += translated_chunk  # Add the translated chunk to the translated text
                break  # Break the while loop if the translation was successful
            except Exception as e:
                print(f"Failed to translate text-english: {chunk}. Error: {str(e)}")
                time.sleep(1)  # Wait for 1 second before retrying

    #translated_text = translated_text.replace("'", "")  # remove apostrophes
    return translated_text

def translate_text_frame(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated = translate_text(run.text)  # Translate the text to English
                run.text = translated  # Update the text with the translated version
            except Exception as e:
                print(f"Failed to translate text-frame: {run.text}. Error: {str(e)}")
                continue

def adjust_text_size(shape):
    while shape.text_frame.text != "":
        try:
            # Try to access the last character of the shape's text
            _ = shape.text_frame.text[-1]
            break  # Break the loop if the last character can be accessed
        except IndexError:
            # If the last character cannot be accessed, the text overflows the shape
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Decrease the font size by 1 point
                    run.font.size = Pt(run.font.size.pt - 1)

def process_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # If the shape is a group, recursively process each shape within the group
        for shape_in_group in shape.shapes:
            process_shape(shape_in_group)
    elif shape.has_text_frame:
        # If the shape has a text frame, translate the text and adjust the text size
        translate_text_frame(shape.text_frame)
        adjust_text_size(shape)
    elif shape.has_table:
        # If the shape is a table, process each cell's text frame and adjust text size
        for row in shape.table.rows:
            for cell in row.cells:
                translate_text_frame(cell.text_frame)
                adjust_text_size(cell)

def translate_pptx(pptx_path):
    pres = Presentation(pptx_path)

    # Process each slide in the presentation
    for slide in pres.slides:
        for shape in slide.shapes:
            process_shape(shape)
            if shape.has_text_frame:
                translate_text_frame(shape.text_frame)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        translate_text_frame(cell.text_frame)

    directory, filename = os.path.split(pptx_path)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".pptx"
    translated_pptx_path = os.path.join(directory, translated_filename)
    pres.save(translated_pptx_path)

    print(f"Translated presentation saved at: {translated_pptx_path}")
    return translated_pptx_path

def convert_doc_to_docx(doc_path):
    # Ensure the path is absolute
    doc_path = os.path.abspath(doc_path)

    # Create the new path by replacing the extension
    new_file_abs = doc_path.replace(".doc", ".docx")

    try:
        # Initialize the Word.Application
        word = win32com.client.Dispatch('Word.Application',pythoncom.CoInitialize())

        # Set the application to be invisible
        word.Visible = False

        # Open the document
        doc = word.Documents.Open(doc_path)

        # Save as a .docx file
        doc.SaveAs(new_file_abs, FileFormat=16)  # 16 represents the wdFormatDocx constant

        # Close the document
        doc.Close()

        # Quit Word
        word.Quit()
    except Exception as e:
        print(f"Failed to convert file: {doc_path}. Error: {str(e)}")
        return doc_path

    return new_file_abs

def translate_docx(doc_path):
    doc = Document(doc_path)

    # Translate text in paragraphs
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated = translate_text(run.text)  # Translate the text to English
                run.text = translated  # Update the text with the translated version
            except Exception as e:
                print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                continue

    # Translate text in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if not run.text:
                            continue  # Skip empty runs
                        try:
                            translated = translate_text(run.text)  # Translate the text to English
                            run.text = translated  # Update the text with the translated version
                        except Exception as e:
                            print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                            continue

    directory, filename = os.path.split(doc_path)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".docx"
    translated_doc_path = os.path.join(directory, translated_filename)
    doc.save(translated_doc_path)

    print(f"Translated document saved at: {translated_doc_path}")
    return translated_doc_path

def create_folder(name, parent_id):
    name = remove_illegal_chars(name)  # Remove illegal characters from the folder name

    # Check if the folder already exists
    response = drive_service.files().list(
        q=f"name='{name}' and '{parent_id}' in parents and mimeType='application/vnd.google-apps.folder'",
        spaces='drive',
        fields='files(id, name)').execute()
    
    if response.get('files'):
        folder_id = response.get('files')[0].get('id')
        print(f"Folder '{name}' already exists in parent folder ID '{parent_id}'")
        return folder_id, False  # Return False indicating the folder already exists

    # Create the folder if it doesn't exist
    file_metadata = {
        'name': name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_id]
    }
    file = drive_service.files().create(body=file_metadata,
                                        fields='id').execute()
    
    print(f"Folder '{name}' created in parent folder ID '{parent_id}'")
    return file.get('id'), True  # Return True indicating the folder is newly created

def delete_file(file_id):
    try:
        drive_service.files().delete(fileId=file_id).execute()
        print(f"File with ID {file_id} has been deleted.")
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

def download_file(file_id, directory_path, file_name, translated_root_id, override=False):
    if not os.path.exists(directory_path):
        os.makedirs(directory_path)  # Create the directory if it doesn't exist

    translated_file_name = translate_file_name(file_name)
    file_path = os.path.join(directory_path, translated_file_name)

    # Check if the file already exists in the translated_root_id directory on Google Drive
    encoded_filename = urllib.parse.quote(translated_file_name)
    response = drive_service.files().list(
        q=f"name='{encoded_filename}' and '{translated_root_id}' in parents",
        fields='files(id, name)').execute()
    
    if response.get('files') and not override:
        print(f"File '{translated_file_name}' already exists in parent folder ID '{translated_root_id}'")
        return None  # Return None if the file already exists
    
    elif response.get('files') and override:
        # Get the ID of the file to delete
        file_to_delete_id = response.get('files')[0]['id']
        delete_file(file_to_delete_id)  # Delete the existing file
    
    # If the file does not exist, download it
    request = drive_service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while done is False:
        status, done = downloader.next_chunk()
    with open(file_path, 'wb') as f:
        f.write(fh.getbuffer())

    return file_path

def upload_file(file_path, parent_folder_id, override=True):
    file_metadata = {
        'name': os.path.basename(file_path),
        'parents': [parent_folder_id]
    }

    # Check if the file already exists in the parent_folder_id directory on Google Drive
    response = drive_service.files().list(
        q=f"name='{file_metadata['name']}' and '{parent_folder_id}' in parents",
        fields='files(id, name)').execute()
    
    if response.get('files') and not override:
        print(f"File '{file_metadata['name']}' already exists in parent folder ID '{parent_folder_id}'")
        return  # Do not upload the file if it already exists
    
    elif response.get('files') and override:
        # Get the ID of the file to delete
        file_to_delete_id = response.get('files')[0]['id']
        delete_file(file_to_delete_id)  # Delete the existing file
    
    media = MediaFileUpload(file_path)
    file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    print(f"File uploaded with ID: {file.get('id')}")

def copy_and_rename_file(file_id, translated_root_id, translated_file_name, override=False):
    # Make a copy of the original file in the new directory
    file_metadata = {
        'name': translated_file_name,
        'parents': [translated_root_id]
    }

    # Check if the file already exists in the translated_root_id directory on Google Drive
    response = drive_service.files().list(
        q=f"name='{translated_file_name}' and '{translated_root_id}' in parents",
        fields='files(id, name)').execute()
    
    if response.get('files') and not override:
        print(f"File '{translated_file_name}' already exists in parent folder ID '{translated_root_id}'")
        return  # Do not copy the file if it already exists
    
    elif response.get('files') and override:
        # Get the ID of the file to delete
        file_to_delete_id = response.get('files')[0]['id']
        delete_file(file_to_delete_id)  # Delete the existing file
    
    copied_file = drive_service.files().copy(
        fileId=file_id,
        body=file_metadata,
        fields='id'
    ).execute()
    print(f"File {translated_file_name} copied and renamed with ID: {copied_file.get('id')}")

def translate_file_name(file_name):
    file_name_without_ext, ext = os.path.splitext(file_name)

    try:
        translated_file_name_without_ext = translate_text(file_name_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {file_name_without_ext}. Error: {str(e)}")
        translated_file_name_without_ext = file_name_without_ext

    translated_file_name_without_ext = remove_illegal_chars(translated_file_name_without_ext)
    translated_file_name = translated_file_name_without_ext + ext
    return translated_file_name

def process_file(item, local_directory_path, translated_root_id, convert_docs, override_docs, convert_slides, override_slides, copy_translate_others, override_others):
    file_name = item['name']
    file_id = item['id']

    # Query if the translated file already exists in the translated directory
    translated_file_name = translate_file_name(file_name)

    # URL encode the filename
    encoded_filename = urllib.parse.quote(translated_file_name)
    response = drive_service.files().list(
        q=f"name='{encoded_filename}' and '{translated_root_id}' in parents",
        fields='files(id, name)').execute()
    if response.get('files'):
        print(f"File '{translated_file_name}' already exists in parent folder ID '{translated_root_id}'")
        return  # Skip this file if the translated version already exists

    # Now download and process the file only if the translated version doesn't exist
    if item['mimeType'] in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'] and convert_docs:
        file_path = download_file(file_id, local_directory_path, file_name, translated_root_id, False)
        if item['mimeType'] == 'application/msword':  # If it is a .doc file
            file_path = convert_doc_to_docx(file_path)
        print(f"Translating file: {translated_file_name}")
        translated_file_path = translate_docx(file_path)
        upload_file(translated_file_path, translated_root_id, override_docs)
        os.remove(translated_file_path)  # Delete the file after uploading
    elif item['mimeType'] == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and convert_slides:
        file_path = download_file(file_id, local_directory_path, file_name, translated_root_id, False)
        print(f"Translating file: {translated_file_name}")
        translated_file_path = translate_pptx(file_path)
        upload_file(translated_file_path, translated_root_id, override_slides)
        os.remove(translated_file_path)  # Delete the file after uploading
    elif copy_translate_others:
        # If the file is not a Word document, copy and rename it without translating the content
        copy_and_rename_file(file_id, translated_root_id, translated_file_name, override_others)

def process_directory(directory_id, translated_root_id, start_file=None, convert_docs=True, override_docs=True, convert_slides=True, override_slides=True, copy_translate_others=True, override_others=True, process_folders=True):
    results = drive_service.files().list(
        q=f"'{directory_id}' in parents",
        fields="files(id, name, mimeType)").execute()

    items = results.get('files', [])

    local_directory_path = os.path.join('./temp_drive_files', directory_id)

    for item in items:
        if item['mimeType'] == 'application/vnd.google-apps.folder':
            subdirectory_id = item['id']
            file_name = translate_text(item['name'])

            # If start_translating is False and (the current item's name matches start_file or start_file is None), set start_translating to True
            global start_translating
            if not start_translating and (not start_file or file_name == start_file):
                start_translating = True

            translated_subdirectory_id, is_new_folder = create_folder(file_name, translated_root_id)
            if is_new_folder or process_folders:  # Only process the subdirectory if it is newly created
                print(f"Processing subdirectory '{item['name']}' with ID '{subdirectory_id}'")
                process_directory(subdirectory_id, translated_subdirectory_id, start_file, convert_docs, override_docs, convert_slides, override_slides, copy_translate_others, override_others)
        else:
            # Skip processing this item if start_translating is False
            if not start_translating:
                continue

            process_file(item, local_directory_path, translated_root_id, convert_docs, override_docs, convert_slides, override_slides, copy_translate_others, override_others)
